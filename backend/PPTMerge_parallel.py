import os
import re
import argparse
import logging
import json
import time
import concurrent.futures
from PIL import Image, ImageDraw, ImageOps
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE, MSO_SHAPE_TYPE, PP_PLACEHOLDER
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import nsdecls
from pptx.oxml import parse_xml
from llm_api import LLMAPIClient
import math 

# Setup logging
logging.basicConfig(level=logging.INFO, format='%(asctime)s - %(levelname)s - %(message)s')

def read_file(file_path):
    with open(file_path, 'r', encoding='utf-8') as f:
        return f.read()

def extract_code_block(text):
    # 1. 尝试匹配完整的 python 代码块
    match = re.search(r'```python\s*(.*?)\s*```', text, re.DOTALL)
    if match: return match.group(1)
    
    # 2. 尝试匹配完整的通用代码块
    match = re.search(r'```\s*(.*?)\s*```', text, re.DOTALL)
    if match: return match.group(1)
    
    # 3. 处理截断的代码块 (缺少闭合的 ```)
    # 这种情况通常发生在 max_tokens 限制导致输出未完成时
    match = re.search(r'```python\s*(.*)', text, re.DOTALL)
    if match: return match.group(1)
    
    match = re.search(r'```\s*(.*)', text, re.DOTALL)
    if match: return match.group(1)
    
    return text

# [新增] 修复截断代码的函数
def try_fix_truncated_code(code):
    """尝试修复截断的 Python 代码 (自动补全括号)"""
    open_brackets = {'(': ')', '[': ']', '{': '}'}
    stack = []
    # 简单的状态机，忽略字符串内的字符
    in_string = False
    string_char = ''
    
    for i, char in enumerate(code):
        if char in ['"', "'"]:
            if not in_string:
                in_string = True
                string_char = char
            elif char == string_char and (i == 0 or code[i-1] != '\\'):
                in_string = False
        
        if not in_string:
            if char in open_brackets:
                stack.append(open_brackets[char])
            elif char in open_brackets.values():
                if stack and stack[-1] == char:
                    stack.pop()
    
    # 补全剩余的括号
    while stack:
        code += stack.pop()
    return code

def natural_sort_key(s):
    """用于对文件名进行自然排序 (1_2 在 1_10 之前)"""
    return [int(text) if text.isdigit() else text.lower()
            for text in re.split('([0-9]+)', s)]

def split_manim_scenes(manim_code):
    scene_pattern = re.compile(r'^class\s+(\w+)\(.*?Scene\):', re.MULTILINE)
    matches = list(scene_pattern.finditer(manim_code))
    scenes = []
    if not matches:
        return [("MainScene", manim_code)]
    
    for i, match in enumerate(matches):
        start_idx = match.start()
        scene_name = match.group(1)
        if i < len(matches) - 1:
            end_idx = matches[i+1].start()
            scene_code = manim_code[start_idx:end_idx]
        else:
            scene_code = manim_code[start_idx:]
        scenes.append((scene_name, scene_code))
    return scenes

# --- 核心辅助函数：坐标转换 ---
# Manim (Default): W=14.22, H=8.0, Center=(0,0)
# PPTX: W=13.33, H=7.5, TopLeft=(0,0)
def manim_to_ppt_coords(manim_x, manim_y):
    """
    [FIXED] Transforms Manim coordinates to PPTX coordinates (Inches).
    Manim X: [-7.11, 7.11] -> PPT Width (13.33 inches)
    Manim Y: [-4.0, 4.0]   -> PPT Height (7.5 inches)
    """
    slide_width_inches = 13.33
    slide_height_inches = 7.5
    manim_width = 14.22
    manim_height = 8.0

    # 1. 将 Manim 坐标（中心为原点）转换到以左上角为原点的系统
    # X: from [-7.11, 7.11] to [0, 14.22]
    # Y: from [4.0, -4.0] to [0, 8.0] (Y轴翻转)
    manim_x_shifted = manim_x + manim_width / 2
    manim_y_shifted = -manim_y + manim_height / 2

    # 2. 按比例缩放到 PPT 的尺寸
    ppt_x = (manim_x_shifted / manim_width) * slide_width_inches
    ppt_y = (manim_y_shifted / manim_height) * slide_height_inches
    
    return Inches(ppt_x), Inches(ppt_y)

def add_arrow_head_func(line):
    ln = line._get_or_add_ln()
    headEnd_xml = '<a:headEnd xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" type="triangle" w="med" len="med"/>'
    headEnd = parse_xml(headEnd_xml)
    ln.append(headEnd)

def add_centered_shape(slide, shape_type, center_x_inch, center_y_inch, width_inch, height_inch):
    """
    以中心点坐标添加形状，自动计算左上角坐标。
    """
    left = center_x_inch - (width_inch / 2)
    top = center_y_inch - (height_inch / 2)
    shape = slide.shapes.add_shape(shape_type, left, top, width_inch, height_inch)
    return shape

def add_centered_textbox(slide, text, center_x_inch, center_y_inch, width_inch, height_inch, word_wrap=True):
    """
    以中心点添加文本框
    """
    try:
        from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    except ImportError:
        pass  # PP_ALIGN, MSO_ANCHOR 已在文件顶部导入

    left = center_x_inch - (width_inch / 2)
    top = center_y_inch - (height_inch / 2)
    box = slide.shapes.add_textbox(left, top, width_inch, height_inch)
    box.text_frame.text = text
    box.text_frame.word_wrap = word_wrap
    box.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    box.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    return box

def fallback_render(slide, scene_name, manim_code, client, exec_globals):
    """
    保底渲染函数。当代码多次修复失败后，使用 LLM 生成仅包含文本的简单幻灯片代码。
    """
    logging.info(f"  -> Triggering LLM-based fallback render for {scene_name}...")
    
    # 1. 清理幻灯片
    for shape in list(slide.shapes):
        sp = shape.element
        sp.getparent().remove(sp)

    # 2. 定义保底模板
    template = r'''
# [Fallback Template]
# 环境中已存在: slide, Inches, Pt, RGBColor, PP_ALIGN, add_centered_textbox 等

# 1. 添加标题
title_y = 0.5
title_width = Inches(12)
left_position = (prs.slide_width - title_width) / 2

title_box = slide.shapes.add_textbox(left_position, Inches(title_y), title_width, Inches(1.0))
tf = title_box.text_frame
tf.word_wrap = False

p = tf.paragraphs[0]
p.text = "SCENE_NAME_PLACEHOLDER" # 请替换为实际标题
p.font.size = Pt(32)
p.font.bold = True
p.font.color.rgb = RGBColor(255, 255, 255)
p.alignment = PP_ALIGN.CENTER

line_width_inch = 9.0
line_height = 0  # 水平线高度为0

line_left = (prs.slide_width - Inches(line_width_inch)) / 2
line_top = Inches(1.1)

line_shape = slide.shapes.add_shape(
    MSO_SHAPE.LINE_INVERSE,
    line_left,
    line_top,
    Inches(line_width_inch),
    line_height
)

line = line_shape.line
line.color.rgb = RGBColor(255, 134, 47)  # Manim Orange (#FF862F)
line.width = Pt(3)
```

# 2. 添加正文文本
# 正文位置: 标题下方
# 正文文本颜色: 白色 (#FFFFFF)
content_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(11.5), Inches(5.5))
content_tf = content_box.text_frame
content_tf.word_wrap = True

# [REPLACE_WITH_EXTRACTED_CONTENT]
# 示例:
# p = content_tf.add_paragraph()
# p.text = "提取的文本行1"
# p.font.size = Pt(24)
# p.space_after = Pt(10)
'''

    # 3. 构造 Prompt
    prompt = f"""
你是一个 Python python-pptx 代码专家。
之前的代码尝试渲染 Manim 场景 "{scene_name}" 失败了。
现在我们需要生成一个**纯文本的保底幻灯片**。

**任务**：
1. 分析提供的 Manim 代码，提取其中的核心文本信息（如 Text, Tex, Title 等内容）。
2. 使用我提供的**代码模板**，生成 Python 代码将这些文本添加到 `slide` 中。
3. **不要**包含任何动画、图形、图片或复杂布局。只保留文本。
4. 代码将在一个已经包含 `slide`, `Inches`, `Pt` 等对象的环境中执行。

**Manim 代码：**
```python
{manim_code}
```

**代码模板：**
```python
{template}
```

**要求**：
- 将模板中的 `SCENE_NAME_PLACEHOLDER` 替换为 "{scene_name}" 或从代码中提取的更合适的标题。
- 提取 Manim 代码中的文本内容，并生成对应的 `content_tf.add_paragraph()` 代码。
- 保持代码简洁，确保可以直接运行。
- 只输出 Python 代码块。
"""

    try:
        # 调用 LLM
        response = client.call_api_with_text(prompt)
        fallback_code = extract_code_block(response)
        
        # 执行生成的代码
        exec(fallback_code, exec_globals)
        logging.info(f"  -> Fallback render successful for {scene_name}")
        
    except Exception as e:
        logging.error(f"  -> Fallback render failed: {e}")
        # 如果 LLM 生成失败，使用最简单的文本框显示错误
        try:
            err_box = slide.shapes.add_textbox(Inches(1), Inches(3), Inches(10), Inches(2))
            err_box.text_frame.text = f"Fallback render failed for {scene_name}.\nError: {e}"
        except:
            pass
    
    footer_box = slide.shapes.add_textbox(Inches(0.5), Inches(6.8), Inches(12), Inches(0.5))
    tf_footer = footer_box.text_frame
    p_footer = tf_footer.paragraphs[0]
    p_footer.text = "注：此页为简化展示，详见原 Manim 场景。"
    p_footer.font.size = Pt(14)
    p_footer.font.color.rgb = RGBColor(128, 128, 128)
    p_footer.alignment = PP_ALIGN.LEFT


def debug_and_execute_code(client, original_code, scene_name, slide, prs, exec_globals):
    """
    执行并调试代码，包含3次重试逻辑。
    返回: (success: bool, final_code: str)
    """
    current_code = original_code
    max_retries = 3

    for attempt in range(max_retries + 1):
        try:
            # 尝试执行代码
            exec(current_code, exec_globals)
            logging.info(f"  -> Success: {scene_name} (Attempt {attempt + 1})")
            return True, current_code
        except Exception as e:
            import traceback
            error_message = f"{e}\n{traceback.format_exc()}"
            logging.warning(f"  -> Attempt {attempt + 1} failed for '{scene_name}': {e}")

            if attempt < max_retries:
                # 如果还没到最大重试次数，请求修复
                logging.info(f"  -> Requesting fix for '{scene_name}'...")
                fix_prompt_path = '/home/TeachMaster/ML/prompt_templates/manim2ppt_fix.txt'
                with open(fix_prompt_path, 'r', encoding='utf-8') as f:
                    fix_prompt_template = f.read()
                fix_prompt = fix_prompt_template.format(
                    current_code=current_code,
                    error_message=error_message
                )
                
                fixed_code_response = client.call_api_with_text(fix_prompt)
                current_code = extract_code_block(fixed_code_response)
                # current_code = try_fix_truncated_code(current_code)
                if current_code.startswith("错误："):
                    logging.error(f"  -> Fix attempt returned error message for '{scene_name}': {current_code}")
                    current_code = original_code  # 回退到原始代码
                    continue
            else:
                # 达到最大重试次数
                logging.error(f"  -> All {max_retries + 1} attempts failed for '{scene_name}'.")
                return False, current_code

    return False, current_code

def generate_pptx_code(client, manim_code, prompt_template):
    prompt = prompt_template.replace("{code}", manim_code)
    response = client.call_api_with_text(prompt)
    return extract_code_block(response)

# --- 新增：样式与品牌辅助函数 ---

def add_background_to_slide(slide, bg_path, prs_width, prs_height):
    """设置背景"""
    if not (bg_path and os.path.exists(bg_path)):
        return
    try:
        try:
            slide.background.fill.user_picture(bg_path)
        except Exception:
            bg = slide.shapes.add_picture(
                bg_path, Inches(0), Inches(0),
                width=prs_width, height=prs_height
            )
            spTree = slide.shapes._spTree
            spTree.remove(bg._element)
            spTree.insert(0, bg._element)
        
        # 确保占位符不遮挡背景
        for shape in slide.shapes:
            try:
                if getattr(shape, "is_placeholder", False):
                    ph_t = shape.placeholder_format.type
                    if ph_t in {
                        PP_PLACEHOLDER.TITLE, PP_PLACEHOLDER.CENTER_TITLE,
                        PP_PLACEHOLDER.SUBTITLE, PP_PLACEHOLDER.BODY, PP_PLACEHOLDER.OBJECT
                    }:
                        shape.fill.background()
            except Exception:
                pass
    except Exception as e:
        logging.warning(f"Failed to add background: {e}")

def add_brand_slide(prs, title, subtitle, teacher_name, teacher_avatar_path, bg_path, left_logo_path, right_logo_path, role="cover"):
    """添加封面或结束页"""
    blank_layout = None
    for lay in prs.slide_layouts:
        if not lay.placeholders:
            blank_layout = lay
            break
    if blank_layout is None:
        blank_layout = prs.slide_layouts[6 if len(prs.slide_layouts) > 6 else 0]
    
    slide = prs.slides.add_slide(blank_layout)

    # Background
    if bg_path and os.path.exists(bg_path):
        add_background_to_slide(slide, bg_path, prs.slide_width, prs.slide_height)

    top_margin = Inches(0.25)
    side_margin = Inches(0.35)
    left_logo_height = Inches(0.6)
    right_logo_height = Inches(0.6)  # Increased the height of the right logo

    # Logos
    if left_logo_path and os.path.exists(left_logo_path):
        try:
            slide.shapes.add_picture(left_logo_path, side_margin, top_margin, height=left_logo_height)
        except Exception: pass
    
    if right_logo_path and os.path.exists(right_logo_path):
        try:
            pic = slide.shapes.add_picture(right_logo_path, Inches(0), Inches(0), height=right_logo_height)
            pic.left = prs.slide_width - pic.width - side_margin
            pic.top = top_margin + Inches(0.05)
        except Exception: pass

    # Title
    title_box = slide.shapes.add_textbox(Inches(1.2), Inches(2.65), Inches(9.5), Inches(2.2))
    # title_box = slide.shapes.add_textbox(Inches(1.2), Inches(2.0), Inches(9.5), Inches(2.2))
    tf = title_box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = title or ""
    p.font.size = Pt(62)
    p.font.bold = True
    p.font.name = "KaiTi"  # Set font to KaiTi
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.LEFT

    # Subtitle (only for cover usually, but logic allows both)
    if subtitle and role == "cover":
        sub_box = slide.shapes.add_textbox(Inches(1.25), Inches(3.85), Inches(9.0), Inches(1.4))
        # sub_box = slide.shapes.add_textbox(Inches(1.25), Inches(3.2), Inches(9.0), Inches(1.4))
        stf = sub_box.text_frame
        stf.clear()
        sp = stf.paragraphs[0]
        sp.text = subtitle
        sp.font.size = Pt(36)
        sp.font.name = "KaiTi"  # Set font to KaiTi
        sp.font.color.rgb = RGBColor(255, 255, 255)
        sp.alignment = PP_ALIGN.LEFT

    # Avatar Circle
    circle_d = Inches(3.6)
    circle_left = prs.slide_width - Inches(5.0)
    circle_top = Inches(1.8)
    
    # Calculate right edge for alignment
    avatar_right_edge = circle_left + circle_d

    if teacher_avatar_path and os.path.exists(teacher_avatar_path):
        try:
            # Use PIL to crop the image to a circle
            with Image.open(teacher_avatar_path) as img:
                img = img.convert("RGBA")
                width, height = img.size
                min_dim = min(width, height)
                
                # Center crop to square
                left = (width - min_dim) / 2
                top = (height - min_dim) / 2
                right = (width + min_dim) / 2
                bottom = (height + min_dim) / 2
                
                img_cropped = img.crop((left, top, right, bottom))
                
                # Create circular mask
                mask = Image.new('L', img_cropped.size, 0)
                draw = ImageDraw.Draw(mask)
                draw.ellipse((0, 0) + img_cropped.size, fill=255)
                
                # Apply mask
                output = ImageOps.fit(img_cropped, mask.size, centering=(0.5, 0.5))
                output.putalpha(mask)
                
                # Save to a temporary file
                temp_avatar_path = os.path.join(os.path.dirname(teacher_avatar_path), "temp_avatar_circle.png")
                output.save(temp_avatar_path)
                
            # Insert the circular image
            slide.shapes.add_picture(temp_avatar_path, circle_left, circle_top, width=circle_d, height=circle_d)
            
            # Clean up
            if os.path.exists(temp_avatar_path):
                os.remove(temp_avatar_path)
                
        except Exception as e:
            logging.warning(f"Failed to process avatar: {e}")
            # Fallback to solid fill if image fails
            oval = slide.shapes.add_shape(MSO_SHAPE.OVAL, circle_left, circle_top, circle_d, circle_d)
            oval.fill.solid()
            oval.fill.fore_color.rgb = RGBColor(255, 255, 255)
            oval.line.fill.background()
    else:
        oval = slide.shapes.add_shape(MSO_SHAPE.OVAL, circle_left, circle_top, circle_d, circle_d)
        oval.fill.solid()
        oval.fill.fore_color.rgb = RGBColor(255, 255, 255)
        oval.line.fill.background()

    # Signature & Link
    sig_text = f"由{teacher_name}和Timo共建" if teacher_name else ""
    
    # Text box dimensions
    tb_width = Inches(6.6)
    tb_height = Inches(1.5)
    
    # Align right edge of text box with right edge of avatar
    tb_left = avatar_right_edge - tb_width
    
    sig_box = slide.shapes.add_textbox(tb_left, prs.slide_height - Inches(2.0),
                                       tb_width, tb_height)
    sg = sig_box.text_frame
    sg.clear()
    
    # Signature Paragraph
    sp = sg.paragraphs[0]
    sp.text = sig_text
    sp.font.size = Pt(24)
    sp.font.name = "KaiTi"  # Set font to KaiTi
    sp.font.color.rgb = RGBColor(255, 255, 255)
    sp.alignment = PP_ALIGN.RIGHT
    
    # Link Paragraph
    link_p = sg.add_paragraph()
    link_p.text = "www.teachmaster.cn"
    link_p.font.size = Pt(18)
    link_p.font.color.rgb = RGBColor(255, 255, 255)
    link_p.alignment = PP_ALIGN.RIGHT
    
    # Add hyperlink
    if link_p.runs:
        r = link_p.runs[0]
        r.hyperlink.address = "http://www.teachmaster.cn"

def put_speech_in_comment(slide, speech_file_path):
    """将讲稿内容添加到幻灯片备注"""
    if not (speech_file_path and os.path.exists(speech_file_path)):
        return
    try:
        with open(speech_file_path, 'r', encoding='utf-8') as f:
            speech_content = f.read().strip()
        
        notes_slide = slide.notes_slide
        text_frame = notes_slide.notes_text_frame
        
        # 添加分隔符
        p = text_frame.add_paragraph()
        p.text = "\n--- 讲稿内容 ---"
        p.font.bold = True
        p.font.size = Pt(20) # Inches(0.28) approx 20pt
        
        # 添加内容
        p = text_frame.add_paragraph()
        p.text = speech_content
        p.font.size = Pt(18) # Inches(0.25) approx 18pt
    except Exception as e:
        logging.warning(f"Failed to add speech comment: {e}")

def add_page_numbers(prs):
    """添加页码 (跳过首页和尾页)"""
    total = len(prs.slides)
    for idx, slide in enumerate(prs.slides, start=1):
        if idx == 1 or idx == total:
            continue
        try:
            box = slide.shapes.add_textbox(
                prs.slide_width - Inches(1.0),
                prs.slide_height - Inches(0.55),
                Inches(0.9),
                Inches(0.4)
            )
            tf = box.text_frame
            tf.clear()
            p = tf.paragraphs[0]
            p.text = str(idx - 1)
            p.font.size = Pt(14)
            p.font.color.rgb = RGBColor(255, 255, 255)
            p.alignment = PP_ALIGN.RIGHT
        except Exception as e:
            logging.warning(f"Failed to add page number: {e}")

def process_scene_task(client, scene_name, scene_code, prompt_template, codes_dir, base_display, scene_index):
    """
    并行任务函数：处理单个场景的代码生成和调试
    返回: (scene_index, success, final_code_or_manim_code)
    """
    try:
        logging.info(f"  [Start] Processing scene: {scene_name}")
        
        # 1. 初始代码生成
        pptx_code = generate_pptx_code(client, scene_code, prompt_template)
        # pptx_code = try_fix_truncated_code(pptx_code)

        # 保存初始生成的代码
        code_filename = f"{base_display}.{scene_index+1}.py"
        with open(os.path.join(codes_dir, code_filename), 'w', encoding='utf-8') as f:
            f.write(pptx_code)
        
        # 2. 准备临时执行环境进行调试
        # 创建一个临时的 Presentation 对象，避免多线程竞争主 prs
        temp_prs = Presentation()
        temp_prs.slide_width = Inches(13.33)
        temp_prs.slide_height = Inches(7.5)
        temp_slide = temp_prs.slides.add_slide(temp_prs.slide_layouts[6])
        
        exec_globals = {
            'slide': temp_slide, 'shapes': temp_slide.shapes, 'prs': temp_prs,
            'Presentation': Presentation, 'Inches': Inches, 'Pt': Pt,
            'RGBColor': RGBColor, 'MSO_CONNECTOR': MSO_CONNECTOR,
            'MSO_SHAPE': MSO_SHAPE, 'PP_ALIGN': PP_ALIGN,
            'nsdecls': nsdecls, 'parse_xml': parse_xml, 'math': math,
            'manim_to_ppt_coords': manim_to_ppt_coords,
            'add_arrow_head': add_arrow_head_func,
            'add_arrow_head_func': add_arrow_head_func,
            'add_centered_shape': add_centered_shape,
            'add_centered_textbox': add_centered_textbox
        }
        
        # 3. 执行并调试代码
        success, final_code = debug_and_execute_code(
            client, pptx_code, scene_name, temp_slide, temp_prs, exec_globals
        )
        
        logging.info(f"  [End] Finished scene: {scene_name} (Success: {success})")
        
        if success:
            return (scene_index, True, final_code)
        else:
            # 失败时返回原始 manim_code 以便主线程进行 fallback_render
            return (scene_index, False, scene_code)

    except Exception as e:
        logging.error(f"  -> Critical error in task {scene_name}: {e}")
        return (scene_index, False, scene_code)

def create_presentation_from_manim(input_path, output_pptx_path, config_path, prompt_template_path,
                                   title=None, subtitle=None, teacher_name=None, teacher_avatar=None,
                                   bg_path=None, left_logo=None, right_logo=None, speech_dir=None,
                                   workers=4):
    
    start_time = time.time()
    logging.info("Starting PPT generation process...")

    # 0. Load Config & Resolve Settings
    config = {}
    if os.path.exists(config_path):
        with open(config_path, 'r', encoding='utf-8') as f:
            import json
            config = json.load(f)
    
    brand = config.get("brand_settings", {})
    
    r_title = title or "Manim Presentation"
    r_subtitle = subtitle or brand.get("subtitle") or "由Timo生成"
    r_teacher = teacher_name or brand.get("teacher_name") or config.get("teacher_name")
    r_avatar = teacher_avatar or brand.get("teacher_avatar_path") or config.get("teacher_avatar_path")
    
    r_bg = bg_path or brand.get("brand_bg_path") or brand.get("global_bg_path") or config.get("brand_bg_path") or config.get("global_bg_path")
    r_left_logo = left_logo or brand.get("left_logo_path")
    r_right_logo = right_logo or brand.get("right_logo_path")

    # 1. Init LLM
    try:
        client = LLMAPIClient(config_path=config_path)
        # 如果 config 中提供了专门的 pptx_settings，优先使用它来覆盖 client 的模型和 base_url
        try:
            pptx_settings = None
            cfg = getattr(client, 'config', None)
            if not cfg and os.path.exists(config_path):
                import json
                with open(config_path, 'r', encoding='utf-8') as _f:
                    cfg = json.load(_f)

            if cfg:
                pptx_settings = cfg.get('pptx_settings', {})

            if pptx_settings:
                logging.info(f"Applying pptx_settings: {pptx_settings}")
                if hasattr(client, 'model'):
                    client.model = pptx_settings.get('model', getattr(client, 'model', None))
                if hasattr(client, 'base_url'):
                    client.base_url = pptx_settings.get('base_url', getattr(client, 'base_url', None))
                
                # [新增] 显式应用 max_tokens
                if hasattr(client, 'max_tokens'):
                    client.max_tokens = pptx_settings.get('max_tokens', getattr(client, 'max_tokens', None))
                # 如果 client 使用 config 字典存储配置
                if hasattr(client, 'config') and isinstance(client.config, dict):
                    client.config['max_tokens'] = pptx_settings.get('max_tokens', client.config.get('max_tokens'))

                # 尝试用 OpenAI SDK 创建底层 client（如果可用），以便使用特定 base_url
                try:
                    from openai import OpenAI
                    api_key = getattr(client, 'api_key', None) or os.environ.get('OPENAI_API_KEY') or (cfg.get('llm_key') if isinstance(cfg, dict) else None)
                    client.client = OpenAI(api_key=api_key, base_url=getattr(client, 'base_url', None))
                except Exception:
                    pass
        except Exception:
            logging.debug('Failed to apply pptx_settings from config, continuing with default client settings')
    except Exception as e:
        logging.error(f"Failed to initialize LLMAPIClient: {e}")
        return

    # 2. Handle Directory
    if os.path.isdir(input_path):
        manim_files = [
            os.path.join(input_path, f) 
            for f in os.listdir(input_path) 
            if f.endswith('.py') and f not in ["0_1.py", "999_1.py"]
        ]
        manim_files.sort(key=lambda x: natural_sort_key(os.path.basename(x)))
    else:
        manim_files = [input_path]

    if not manim_files:
        logging.error("No python files found.")
        return

    prompt_template = read_file(prompt_template_path)

    # 3. Init PPT
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)
    
    output_dir = os.path.dirname(output_pptx_path)
    codes_dir = os.path.join(output_dir, "pptx_codes")
    os.makedirs(codes_dir, exist_ok=True)

    # 4. Add Cover Slide
    add_brand_slide(prs, r_title, r_subtitle, r_teacher, r_avatar, r_bg, r_left_logo, r_right_logo, role="cover")

    # 5. Parallel Processing
    tasks = []
    global_scene_index = 0
    
    # 收集所有任务
    for manim_file_path in manim_files:
        logging.info(f"Collecting scenes from: {manim_file_path}")
        manim_code = read_file(manim_file_path)
        scenes = split_manim_scenes(manim_code)
        base_name = os.path.splitext(os.path.basename(manim_file_path))[0]
        base_display = base_name.replace("_", ".")

        for i, (scene_name, scene_code) in enumerate(scenes):
            # 构造任务元组
            task_info = {
                'client': client,
                'scene_name': scene_name,
                'scene_code': scene_code,
                'prompt_template': prompt_template,
                'codes_dir': codes_dir,
                'base_display': base_display,
                'scene_index': global_scene_index,
                'base_name': base_name # 用于查找讲稿
            }
            tasks.append(task_info)
            global_scene_index += 1

    logging.info(f"Total scenes to process: {len(tasks)}")
    
    # 并行执行
    results = [None] * len(tasks)
    max_workers = min(len(tasks), workers) # 限制最大并发数
    if max_workers < 1: max_workers = 1

    with concurrent.futures.ThreadPoolExecutor(max_workers=max_workers) as executor:
        future_to_index = {
            executor.submit(
                process_scene_task, 
                t['client'], t['scene_name'], t['scene_code'], 
                t['prompt_template'], t['codes_dir'], t['base_display'], t['scene_index']
            ): t['scene_index']
            for t in tasks
        }
        
        for future in concurrent.futures.as_completed(future_to_index):
            idx = future_to_index[future]
            try:
                result = future.result()
                results[idx] = result # (index, success, code_or_manim)
            except Exception as exc:
                logging.error(f"Task {idx} generated an exception: {exc}")
                # 标记为失败
                results[idx] = (idx, False, tasks[idx]['scene_code'])

    # 6. Sequential Assembly
    logging.info("Assembling presentation...")
    
    for i, res in enumerate(results):
        if res is None: continue
        
        idx, success, content = res
        task = tasks[idx]
        scene_name = task['scene_name']
        base_name = task['base_name']
        
        logging.info(f"  -> Adding slide for {scene_name}")
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        
        # 准备执行环境 (供成功执行和 fallback 使用)
        exec_globals = {
            'slide': slide, 'shapes': slide.shapes, 'prs': prs,
            'Presentation': Presentation, 'Inches': Inches, 'Pt': Pt,
            'RGBColor': RGBColor, 'MSO_CONNECTOR': MSO_CONNECTOR,
            'MSO_SHAPE': MSO_SHAPE, 'PP_ALIGN': PP_ALIGN,
            'nsdecls': nsdecls, 'parse_xml': parse_xml, 'math': math,
            'manim_to_ppt_coords': manim_to_ppt_coords,
            'add_arrow_head': add_arrow_head_func,
            'add_arrow_head_func': add_arrow_head_func,
            'add_centered_shape': add_centered_shape,
            'add_centered_textbox': add_centered_textbox
        }

        if success:
            # 执行成功的代码
            final_code = content
            try:
                exec(final_code, exec_globals)
            except Exception as e:
                logging.error(f"Error executing final code for {scene_name}: {e}")
                fallback_render(slide, scene_name, task['scene_code'], client, exec_globals)
        else:
            # 失败，使用 fallback
            fallback_render(slide, scene_name, content, client, exec_globals)

        # 应用背景和讲稿
        add_background_to_slide(slide, r_bg, prs.slide_width, prs.slide_height)
        if speech_dir:
            speech_file = os.path.join(speech_dir, f"{base_name}.txt")
            put_speech_in_comment(slide, speech_file)

    # 7. Add End Slide
    add_brand_slide(prs, r_title, None, r_teacher, r_avatar, r_bg, r_left_logo, r_right_logo, role="end")

    # 8. Add Page Numbers
    add_page_numbers(prs)

    prs.save(output_pptx_path)
    
    end_time = time.time()
    duration = end_time - start_time
    logging.info(f"Saved presentation to {output_pptx_path}")
    logging.info(f"Total execution time: {duration:.2f} seconds")

if __name__ == "__main__":
    parser = argparse.ArgumentParser(description="Convert Manim file to PPTX.")
    parser.add_argument("input_file", help="Path to input Manim .py file or directory")
    parser.add_argument("output_file", help="Path to output .pptx file")
    
    # Default paths
    script_dir = os.path.dirname(os.path.abspath(__file__))
    default_config = os.path.join(script_dir, "config.json")
    default_prompt = os.path.join(script_dir, "prompt_templates", "manim2pptx.txt")
    
    parser.add_argument("--config", default=default_config, help="Path to config.json")
    parser.add_argument("--prompt", default=default_prompt, help="Path to prompt template")
    
    # Brand args
    parser.add_argument("--title", help="Course Title")
    parser.add_argument("--subtitle", help="Subtitle")
    parser.add_argument("--teacher", help="Teacher Name")
    parser.add_argument("--avatar", default="/home/TeachMasterAppV2/backend/ppt_templates/TeachMaster.png", help="Path to teacher avatar")
    parser.add_argument("--bg", default="/home/TeachMasterAppV2/backend/background_default.png", help="Path to background image")
    parser.add_argument("--left_logo", default="/home/TeachMasterAppV2/backend/ppt_templates/sjtupic.png", help="Path to left logo")
    parser.add_argument("--right_logo", default="/home/TeachMasterAppV2/backend/ppt_templates/TeachMaster.png", help="Path to right logo")
    parser.add_argument("--speech_dir", help="Path to speech text directory")
    parser.add_argument("--workers", type=int, default=2, help="Parallel workers for debugging (default: 2)")

    args = parser.parse_args()
    
    create_presentation_from_manim(
        args.input_file, args.output_file, args.config, args.prompt,
        title=args.title, subtitle=args.subtitle, teacher_name=args.teacher,
        teacher_avatar=args.avatar
        , bg_path=args.bg, left_logo=args.left_logo, right_logo=args.right_logo,
        speech_dir=args.speech_dir, workers=args.workers
    )